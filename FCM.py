# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- 讓簡報預設用「微軟正黑體」以免中文變方塊 ---
    def set_font(run, size=Pt(24), bold=False):
        run.font.name = 'Microsoft JhengHei'
        run.font.size = size
        run.font.bold = bold
        run.font.language_id = 0x0404   # 中文語系

    # Helper：一般大綱式投影片
    def add_slide(title, content_points):
        slide_layout = prs.slide_layouts[1]   # 標題與內容
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, txt in enumerate(content_points):
            p = body.add_paragraph()
            p.text = txt
            p.level = 0
            p.space_before = Pt(6) if i == 0 else Pt(12)
            for run in p.runs:
                set_font(run, size=Pt(22))

    # Helper：章節首頁
    def add_section_header(title, subtitle):
        slide_layout = prs.slide_layouts[0]   # 標題投影片
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        # 一併把標題字型改掉
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            set_font(run, size=Pt(40), bold=True)
        for run in slide.placeholders[1].text_frame.paragraphs[0].runs:
            set_font(run, size=Pt(28))

    # ---- 以下內容與你原來的相同，僅刪掉 [cite_start] ----
    add_section_header("管理組織的生存環境：生態與體制理論",
                       "Chapter 06: 組織如何在叢林與社會中求生？")

    add_slide("組織生存的雙重戰場", [
        "組織不只是內部管理，更活在兩個外部環境中：",
        "1. 生物叢林 (組織生態理論)：",
        "   - 關注資源競爭、環境負載力。",
        "   - 核心邏輯：適者生存，位置決定命運。",
        "2. 人類社會 (體制理論)：",
        "   - 關注社會規範、價值觀、正當性。",
        "   - 核心邏輯：合乎禮教者生存，尋求社會認同。"
    ])

    add_section_header("第一部分：組織生態理論",
                       "Organizational Ecology\n叢林法則：資源與利基的戰爭")

    add_slide("環境負載力與族群密度", [
        "環境負載力 (Carrying Capacity)：",
        "   - 每個產業(池塘)能養活的組織(魚)數量是有上限的。",
        "   - 取決於市場規模、資源多寡。",
        "族群密度 (Population Density)：",
        "   - 競爭者的數量。",
        "   - 初期密度增加有助於合法性(把餅做大)，後期則導致惡性競爭。"
    ])

    add_slide("【實務個案】台灣醫療機構的消長 [cite: 20-30]", [
        "現象 (2004-2014)：",
        "   - 醫療總需求上升 (人口老化、慢性病支出成長 54.1%)。",
        "   - 但「醫院」家數減少，「診所」大幅增加 (尤其包含中醫診所)。",
        "教授解析 (環境負載力運作)：",
        "   - 資源限制：護理人力短缺、健保總額支付制度 (點值稀釋)。",
        "   - 大型醫院生存空間被壓縮，靈活的小型診所適應力更好。"
    ])

    add_slide("結構慣性：為什麼大象難跳舞？ [cite: 99-109]", [
        "結構慣性 (Structural Inertia)：組織不易改變核心形式的傾向。",
        "慣性的雙面刃：",
        "   - 優點 (生存基礎)：建立可靠性 (Reliability) 與 課責性 (Accountability)，讓客戶信任。",
        "   - 缺點 (老化風險)：當環境劇變時，SOP 變成阻礙，導致反應遲鈍。",
        "新創的風險 (Liability of Newness)：",
        "   - 新組織雖靈活，但缺乏穩定系統與人脈，死亡率高。"
    ])

    add_slide("生存策略：瑞士刀 vs. 手術刀 [cite: 136-146]", [
        "廣泛經營者 (Generalist)：",
        "   - 像瑞士刀，功能多、市場大。",
        "   - 資源寬裕，分散風險，但面臨激烈競爭。",
        "專精經營者 (Specialist)：",
        "   - 像手術刀，專注特定利基 (Niche)。",
        "   - 資源少，但深耕小市場，避開大公司砲火。",
        "利基寬度理論：",
        "   - 環境變動大且長時，廣泛者較佳；環境穩定或變動短暫，專精者勝。"
    ])

    add_slide("【實務個案】台灣航空業的資源分割 [cite: 251-262]", [
        "資源分割理論 (Resource Partitioning)：大廠吃主流，小廠吃邊緣。",
        "立榮航空 (廣泛經營者)：",
        "   - 經營 14 條航線 (本島+離島)。",
        "   - 面臨華信、遠東的高度競爭 (重疊度高)。",
        "德安航空 (專精經營者)：",
        "   - 專注台東-蘭嶼、綠島等離島航線。",
        "   - 雖然市場小，但「獨佔」資源 (競爭者重疊為 0)。",
        "啟示：小公司不要去紅海硬碰硬，要找大魚游不進去的縫隙。"
    ])

    add_slide("紅皇后效應 (Red Queen Effect) [cite: 167-174]", [
        "「在這裡，你必須盡全力奔跑，才能保持在原地。」——《愛麗絲鏡中奇遇》",
        "定義：",
        "   - 當競爭對手變強，你進步只是為了「不被淘汰」，而非拉開差距。",
        "   - 案例：UPS vs. FedEx 的軍備競賽。",
        "啟示：",
        "   - 遠離競爭或停止學習的組織，會立刻被環境淘汰。"
    ])

    add_slide("第二部分：體制理論", [
        "Institutional Theory",
        "社會規則：不只搶錢，還要搶「面子」與「認同」"
    ])

    add_slide("體制的三大支柱與同形化 [cite: 185-215]", [
        "為何同一產業的公司長得越來越像 (Isomorphism)？",
        "1. 法規性 (Regulative) - 「不得不做」：",
        "   - 政府法令、強制規定 (例：Sony 要求供應商交環保報告書)。",
        "2. 規範性 (Normative) - 「應該要做」：",
        "   - 道德、職業標準 (例：會計師事務所的標準化作業)。",
        "3. 認知性 (Cognitive) - 「大家都這樣做」：",
        "   - 模仿成功者最安全 (例：腳踏車廠模仿豐田式管理)。"
    ])

    add_slide("正當性 (Legitimacy)：組織的社會通行證 [cite: 273-280]", [
        "組織遵循體制，是為了換取「正當性」。",
        "三種正當性衡量：",
        "   1. 實質利益正當性：對利害關係人有用 (如產品好用)。",
        "   2. 道德正當性：符合社會道德 (如環保、不血汗)。",
        "   3. 認知正當性：合乎常理與社會期待 (如銀行建築要氣派)。",
        "後果：",
        "   - 失去正當性，資源獲取將困難，甚至面臨抵制。"
    ])

    add_slide("【實務與理論對話】Facebook 的體制困境 [cite: 328-342]", [
        "Facebook 面臨的挑戰：",
        "   - 法規性體制：早期隨意傳輸個資，違反歐盟法規 (被迫調整或退出)。",
        "   - 規範性體制：假新聞、演算法操弄選舉，違背普世價值 (新聞自由與道德)。",
        "後果：",
        "   - 道德正當性受損，企業形象下跌，面臨更嚴格監管。",
        "建議：",
        "   - 必須從「操控/忽視」轉向「遵從」或更積極的體制管理。"
    ])

    add_slide("教授錦囊：面對體制壓力的四種策略 [cite: 299-312]", [
        "1. 遵從 (Acquiesce)：",
        "   - 乖乖照做，滿足利害關係人 (例：取得 ISO 認證)。",
        "2. 選擇 (Select)：",
        "   - 挑選能接納自己的環境 (例：產品不環保，就賣往法規寬鬆地區)。",
        "3. 操控 (Manipulate)：",
        "   - 改變外界觀感或規則 (例：透過廣告、遊說改變標準)。",
        "4. 緩衝/隱瞞 (Buffer/Conceal)：",
        "   - 表面符合，實際脫鉤 (例：形式上的品管，實際未執行 - 風險高)。"
    ])

    add_slide("結語：組織生存公式", [
        "組織生存 = (找到適合的生態利基) + (取得社會的正當性)",
        "課後思考：",
        "   1. 你的組織是廣泛經營者(大鯨魚)還是專精經營者(小蝦米)？",
        "   2. 你的組織是否具備足夠的「正當性」？社會大眾信任你嗎？",
        "",
        "Next Step: 請分組討論，選擇一個新興產業分析其體制環境。"
    ])

    # 存檔
    file_name = "Ch06_組織生態與體制理論_授課講義.pptx"
    prs.save(file_name)
    print(f"Presentation saved as {file_name}")

if __name__ == "__main__":
    create_presentation()
